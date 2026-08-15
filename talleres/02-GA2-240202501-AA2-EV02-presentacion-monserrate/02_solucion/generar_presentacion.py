from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor as DocxRGBColor


SOLUTION_DIR = Path(__file__).resolve().parent
WORKSHOP_DIR = SOLUTION_DIR.parent
DELIVERY_DIR = WORKSHOP_DIR / "03_entrega"
IMAGE_DIR = SOLUTION_DIR / "recursos" / "imagenes"
DELIVERY_DIR.mkdir(parents=True, exist_ok=True)

PPTX_PATH = DELIVERY_DIR / "GA2-240202501-AA2-EV02_Presentacion_Monserrate.pptx"
GUIDE_PATH = DELIVERY_DIR / "GA2-240202501-AA2-EV02_Guion_Oral.docx"
MARKDOWN_PATH = SOLUTION_DIR / "GUION_PRESENTACION.md"

BASILICA_IMAGE = IMAGE_DIR / "basilica_monserrate.jpg"
CABLE_CAR_IMAGE = IMAGE_DIR / "teleferico_monserrate.jpg"
PANORAMA_IMAGE = IMAGE_DIR / "panorama_bogota.jpg"

GREEN = "39A900"
DARK_GREEN = "123B2A"
DEEP_GREEN = "0B2A1D"
SKY = "83C9DD"
CREAM = "F4F1E8"
PALE_GREEN = "EAF4E5"
TERRACOTTA = "C75A38"
INK = "17231C"
MUTED = "617067"
WHITE = "FFFFFF"
LIGHT_GRAY = "E3E8E3"


SLIDES = [
    {
        "title": "Monserrate: Bogotá from the Sky",
        "time": "0:00–0:35",
        "screen": ["Jhon Steven Alvarez Ruiz"],
        "script": (
            "Good morning. My name is Jhon Steven Alvarez Ruiz. Today I am going to present "
            "Monserrate, one of the most famous tourist places in Bogotá. I chose this destination "
            "because it brings together nature, religion, culture, food, and an excellent view of "
            "the city. I will explain its location, its main elements, its activities, my opinion, "
            "and some recommendations for visitors."
        ),
        "focus": "famous · tourist · destination · together · recommendations",
    },
    {
        "title": "Location and Access",
        "time": "0:35–1:10",
        "screen": [
            "Eastern Hills of Bogotá",
            "Santa Fe locality",
            "3,152 meters above sea level",
            "Cable car · Funicular · Walking trail",
        ],
        "script": (
            "Monserrate is in the Eastern Hills of Bogotá, near the historic city center and in "
            "the locality of Santa Fe. Its summit is 3,152 meters above sea level. Visitors can "
            "reach the top by cable car, funicular, or pedestrian trail. Each option offers a "
            "different experience, but walking is the most demanding because the hill is high "
            "and steep."
        ),
        "focus": "Eastern · historic · locality · summit · pedestrian · demanding",
    },
    {
        "title": "Main Elements",
        "time": "1:10–1:45",
        "screen": [
            "Basilica of the Fallen Lord",
            "Way of the Cross",
            "Panoramic viewpoints",
            "Cable car and funicular",
            "High Andean forest",
            "Restaurants and local food",
        ],
        "script": (
            "At the summit, there is the Basilica of the Fallen Lord of Monserrate. Inside the "
            "church, visitors can see an important religious image. Around the basilica, there is "
            "a Way of the Cross, and there are panoramic viewpoints, transport stations, "
            "restaurants, and cafés. The area also has High Andean forest. For this reason, "
            "Monserrate has religious, cultural, urban, and natural elements."
        ),
        "focus": "summit · basilica · religious · panoramic · Andean · cultural",
    },
    {
        "title": "Description of the Place",
        "time": "1:45–2:20",
        "screen": [
            "High · Green · Steep",
            "Cool and changeable weather",
            "A white basilica at the top",
            "A wide view of Bogotá",
            "Nature · Culture · Religion · City",
        ],
        "script": (
            "Monserrate is a high, green, and impressive hill. Its paths are steep, and the "
            "weather can be cool and changeable. The white basilica stands out against the green "
            "mountain. From the viewpoint, Bogotá looks enormous. Visitors can see many "
            "neighborhoods, buildings, and roads. The outdoor area can be lively, but the basilica "
            "offers a calm and respectful atmosphere."
        ),
        "focus": "impressive · changeable · viewpoint · enormous · neighborhoods · atmosphere",
    },
    {
        "title": "Activities",
        "time": "2:20–3:00",
        "screen": [
            "Take panoramic photographs",
            "Visit the sanctuary",
            "Ride the cable car or funicular",
            "Hike 2.35 km and 1,605 steps",
            "Try Colombian food and coffee",
            "Observe birds on the Paramuno Trail",
        ],
        "script": (
            "People can do many activities at Monserrate. They can ride the cable car or "
            "funicular, take photographs, admire the city, and visit the sanctuary. They can also "
            "try Colombian food and drink coffee. People in good physical condition can hike the "
            "pedestrian trail. It is 2.35 kilometers long and has 1,605 steps. Nature lovers can "
            "enjoy a birdwatching experience on the Paramuno Trail."
        ),
        "focus": "photographs · sanctuary · physical · kilometers · birdwatching · Paramuno",
    },
    {
        "title": "My Opinion",
        "time": "3:00–3:35",
        "screen": [
            "One of Bogotá’s most complete attractions",
            "My favorite feature: the panoramic view",
            "Nature + City + Tradition",
            "Recommended for families, friends, and tourists",
        ],
        "script": (
            "In my opinion, Monserrate is one of the most complete attractions in Bogotá. The "
            "panoramic view is my favorite feature because it helps us understand the size of the "
            "capital. I also like the combination of a modern city, an old religious tradition, "
            "and nature. It is a good destination for families, friends, and tourists. I would "
            "recommend it to anyone who wants to see another side of Bogotá."
        ),
        "focus": "opinion · complete · attractions · favorite · combination · recommend",
    },
    {
        "title": "Recommendations",
        "time": "3:35–4:20",
        "screen": [
            "Check the official opening hours",
            "Wear comfortable shoes",
            "Bring water, sunscreen, and a rain jacket",
            "Walk slowly because of the altitude",
            "Use official routes and follow instructions",
            "Keep the place clean and respect worshipers",
        ],
        "script": (
            "Before going, visitors should check the official opening hours because they can "
            "change. They should wear comfortable shoes and bring water, sunscreen, and a light "
            "rain jacket. At this altitude, it is important to walk slowly and rest if you feel "
            "tired or dizzy. Visitors must stay on the official route and follow the staff’s "
            "instructions. If the hike is too difficult, they can use the funicular or cable car. "
            "Finally, everyone should keep the place clean and respect people who are praying."
        ),
        "focus": "comfortable · sunscreen · altitude · dizzy · route · worshipers",
    },
    {
        "title": "Conclusion",
        "time": "4:20–4:40",
        "screen": [
            "Monserrate connects Bogotá with nature, culture, and tradition.",
            "Thank you for watching.",
        ],
        "script": (
            "In conclusion, Monserrate is more than a viewpoint. It is a place where visitors can "
            "learn about Bogotá, enjoy nature, and understand a local tradition. I hope this "
            "presentation helps you discover this special destination. Thank you for watching."
        ),
        "focus": "conclusion · viewpoint · tradition · presentation · discover",
    },
]


SOURCES = [
    (
        "Visit Bogotá — Monserrate Hill",
        "https://visitbogota.co/en/what-to-do-in-bogota/nature/monserrate-hill",
    ),
    (
        "Cerro de Monserrate — What to do",
        "https://monserrate.co/en/about-monserrate/what-to-do/",
    ),
    (
        "Cerro de Monserrate — Prepare your visit",
        "https://monserrate.co/en/plan-your-visit/",
    ),
    (
        "IDRD — Sendero de Monserrate",
        "https://www.idrd.gov.co/parques-y-escenarios/sendero-de-monserrate",
    ),
]


IMAGE_CREDITS = [
    (
        "Basilica photograph: Jimmy Gómez N — CC BY-SA 3.0",
        "https://commons.wikimedia.org/wiki/File:Monserrate_Bogota_-_panoramio.jpg",
    ),
    (
        "Cable car photograph: Javadumper — CC BY-SA 3.0",
        "https://commons.wikimedia.org/wiki/File:Teleferico_Monserrate.jpg",
    ),
    (
        "Bogotá panorama: diego_cue — CC BY-SA 3.0",
        "https://commons.wikimedia.org/wiki/File:Vistas_desde_Monserrate_-_Bogot%C3%A1_-_panoramio.jpg",
    ),
]


PRONUNCIATION = [
    ("Eastern", "/ˈiːstərn/", "ÍS-tern"),
    ("hills", "/hɪlz/", "jilz, con h suave aspirada"),
    ("summit", "/ˈsʌmɪt/", "SÁ-mit"),
    ("meters", "/ˈmiːtərz/", "MÍ-terz"),
    ("cable car", "/ˈkeɪbəl kɑːr/", "KÉI-bol kar"),
    ("funicular", "/fjuːˈnɪkjələr/", "fiu-NÍ-kiu-ler"),
    ("pedestrian", "/pəˈdestriən/", "pe-DÉS-tri-an"),
    ("basilica", "/bəˈsɪlɪkə/", "be-SÍ-li-ka"),
    ("religious", "/rɪˈlɪdʒəs/", "ri-LÍ-yas"),
    ("panoramic", "/ˌpænəˈræmɪk/", "pa-na-RÁ-mik"),
    ("Andean", "/ˈændiən/", "ÁN-di-an"),
    ("viewpoint", "/ˈvjuːpɔɪnt/", "VIÚ-point"),
    ("neighborhoods", "/ˈneɪbərhʊdz/", "NÉI-bor-hudz"),
    ("sanctuary", "/ˈsæŋktʃuˌeri/", "SÁNK-chu-e-ri"),
    ("birdwatching", "/ˈbɜːrdˌwɒtʃɪŋ/", "BERD-uó-ching"),
    ("attractions", "/əˈtrækʃənz/", "a-TRÁK-shonz"),
    ("comfortable", "/ˈkʌmftəbəl/", "KÁMF-ta-bol"),
    ("altitude", "/ˈæltɪtuːd/", "ÁL-ti-tud"),
    ("worshipers", "/ˈwɜːrʃəpərz/", "UÉR-sha-perz"),
    ("recommend", "/ˌrekəˈmend/", "re-ka-MÉND"),
]


def pptx_color(value):
    return PptxRGBColor.from_string(value)


def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = pptx_color(color)


def add_rect(slide, x, y, width, height, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = pptx_color(fill)
    if line:
        shape.line.color.rgb = pptx_color(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    width,
    height,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font="Lato",
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = pptx_color(color)
    return box


def add_picture_cover(slide, image_path, x, y, width, height):
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    frame_ratio = width / height
    picture = slide.shapes.add_picture(
        str(image_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height)
    )
    if image_ratio > frame_ratio:
        crop = (1 - frame_ratio / image_ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < frame_ratio:
        crop = (1 - image_ratio / frame_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_footer(slide, number, dark=False):
    color = WHITE if dark else MUTED
    add_text(
        slide,
        "GA2-240202501-AA2-EV02  ·  JHON STEVEN ALVAREZ RUIZ",
        0.72,
        7.13,
        9.4,
        0.2,
        size=6.8,
        color=color,
        bold=True,
    )
    add_text(
        slide,
        f"{number:02d} / 08",
        11.7,
        7.1,
        0.9,
        0.22,
        size=7.2,
        color=color,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_header(slide, section, title, dark=False):
    primary = WHITE if dark else INK
    accent = SKY if dark else GREEN
    add_text(slide, section, 0.72, 0.42, 2.4, 0.25, size=8.5, color=accent, bold=True)
    add_text(slide, title, 0.72, 0.7, 10.8, 0.55, size=24, color=primary, bold=True)
    add_rect(slide, 0.72, 1.2, 11.9, 0.035, accent)


def add_pill(slide, text, x, y, width, fill=PALE_GREEN, color=DARK_GREEN, size=10.5):
    pill = add_rect(slide, x, y, width, 0.42, fill, radius=True)
    pill.adjustments[0] = 0.5
    add_text(
        slide,
        text,
        x + 0.07,
        y + 0.035,
        width - 0.14,
        0.31,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def add_numbered_card(slide, number, text, x, y, width, height, fill=WHITE, dark=False):
    add_rect(slide, x, y, width, height, fill, line=LIGHT_GRAY if fill == WHITE else None, radius=True)
    number_fill = GREEN if not dark else SKY
    number_color = WHITE if not dark else DEEP_GREEN
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.2), Inches(0.42), Inches(0.42))
    circle.fill.solid()
    circle.fill.fore_color.rgb = pptx_color(number_fill)
    circle.line.fill.background()
    add_text(
        slide,
        f"{number:02d}",
        x + 0.19,
        y + 0.235,
        0.4,
        0.28,
        size=8.8,
        color=number_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    add_text(
        slide,
        text,
        x + 0.75,
        y + 0.18,
        width - 0.92,
        height - 0.33,
        size=11.4,
        color=WHITE if dark else INK,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_hyperlink_line(slide, label, url, x, y, width, size=7.2, color=DARK_GREEN):
    box = add_text(slide, label, x, y, width, 0.21, size=size, color=color)
    run = box.text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = url
    return box


def build_presentation():
    for image_path in (BASILICA_IMAGE, CABLE_CAR_IMAGE, PANORAMA_IMAGE):
        if not image_path.exists():
            raise FileNotFoundError(f"Missing image: {image_path}")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "Monserrate: Bogotá from the Sky"
    presentation.core_properties.subject = "GA2-240202501-AA2-EV02 — Tourist site presentation"
    presentation.core_properties.author = "Jhon Steven Alvarez Ruiz"
    presentation.core_properties.last_modified_by = "Jhon Steven Alvarez Ruiz"
    presentation.core_properties.keywords = "SENA, English, Monserrate, Bogotá, tourist site"

    blank = presentation.slide_layouts[6]

    # Slide 1 — Cover.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, DEEP_GREEN)
    add_picture_cover(slide, BASILICA_IMAGE, 5.35, 0, 7.98, 7.5)
    add_rect(slide, 5.22, 0, 0.13, 7.5, GREEN)
    add_text(slide, "SENA  ·  ENGLISH PRESENTATION", 0.72, 0.48, 4.1, 0.25, size=8.5, color=SKY, bold=True)
    add_text(slide, "MONSERRATE", 0.72, 1.18, 4.2, 0.7, size=31, color=WHITE, bold=True)
    add_text(slide, "Bogotá from the Sky", 0.72, 1.92, 4.1, 0.5, size=18, color=SKY, bold=True)
    add_text(
        slide,
        "A tourist site where nature, culture, religion, and city life meet.",
        0.72,
        2.6,
        3.95,
        1.0,
        size=12.5,
        color=WHITE,
    )
    add_rect(slide, 0.72, 4.35, 3.92, 0.85, CREAM, radius=True)
    add_text(slide, "Jhon Steven Alvarez Ruiz", 0.94, 4.62, 3.45, 0.3, size=12, color=INK, bold=True)
    add_text(slide, "GA2-240202501-AA2-EV02", 0.72, 6.78, 4.1, 0.25, size=8, color=LIGHT_GRAY, bold=True)
    add_rect(slide, 8.22, 6.6, 5.11, 0.9, GREEN)
    add_text(
        slide,
        "BOGOTÁ FROM THE SKY",
        8.5,
        6.84,
        4.5,
        0.32,
        size=11.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Slide 2 — Location.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, CREAM)
    add_header(slide, "01  ·  WHERE IS IT?", "Location and Access")
    add_text(slide, "EASTERN HILLS", 0.78, 1.6, 4.25, 0.45, size=23, color=DARK_GREEN, bold=True)
    add_text(slide, "Santa Fe · Bogotá · Colombia", 0.78, 2.1, 4.1, 0.35, size=14.5, color=MUTED, bold=True)
    add_rect(slide, 0.78, 2.75, 3.95, 1.35, DARK_GREEN, radius=True)
    add_text(slide, "3,152 m", 1.04, 2.98, 2.1, 0.48, size=27, color=WHITE, bold=True)
    add_text(slide, "ABOVE SEA LEVEL", 1.05, 3.53, 2.75, 0.25, size=8.2, color=SKY, bold=True)
    add_pill(slide, "CABLE CAR", 0.78, 4.55, 1.28, fill=GREEN, color=WHITE, size=9.2)
    add_pill(slide, "FUNICULAR", 2.18, 4.55, 1.28, fill=PALE_GREEN, color=DARK_GREEN, size=9.2)
    add_pill(slide, "WALKING TRAIL", 3.58, 4.55, 1.42, fill=TERRACOTTA, color=WHITE, size=8.8)
    add_text(
        slide,
        "Near Bogotá’s historic city center",
        0.78,
        5.4,
        4.15,
        0.65,
        size=13.5,
        color=INK,
        bold=True,
    )
    add_picture_cover(slide, PANORAMA_IMAGE, 5.35, 1.43, 7.27, 5.38)
    add_rect(slide, 5.35, 6.42, 7.27, 0.39, DARK_GREEN)
    add_text(slide, "A VIEW OVER THE CAPITAL", 5.65, 6.51, 6.7, 0.2, size=8.2, color=WHITE, bold=True)
    add_footer(slide, 2)

    # Slide 3 — Elements.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, DEEP_GREEN)
    add_header(slide, "02  ·  WHAT CAN VISITORS FIND?", "Main Elements", dark=True)
    add_picture_cover(slide, BASILICA_IMAGE, 0.72, 1.5, 4.15, 2.25)
    add_picture_cover(slide, CABLE_CAR_IMAGE, 0.72, 4.02, 4.15, 2.25)
    elements = SLIDES[2]["screen"]
    for index, element in enumerate(elements):
        column = index % 2
        row = index // 2
        add_numbered_card(
            slide,
            index + 1,
            element,
            5.25 + column * 3.75,
            1.55 + row * 1.53,
            3.4,
            1.12,
            fill=DARK_GREEN,
            dark=True,
        )
    add_text(
        slide,
        "RELIGIOUS  ·  CULTURAL  ·  URBAN  ·  NATURAL",
        5.3,
        6.45,
        7.15,
        0.3,
        size=9.5,
        color=SKY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 3, dark=True)

    # Slide 4 — Description.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, WHITE)
    add_header(slide, "03  ·  WHAT IS IT LIKE?", "Description of the Place")
    add_picture_cover(slide, PANORAMA_IMAGE, 0.72, 1.48, 7.45, 4.95)
    adjectives = [
        ("HIGH · GREEN · STEEP", "the mountain"),
        ("COOL · CHANGEABLE", "the weather"),
        ("WIDE · IMPRESSIVE", "the city view"),
        ("CALM · RESPECTFUL", "inside the basilica"),
    ]
    for index, (adjective, context) in enumerate(adjectives):
        y = 1.6 + index * 1.18
        add_rect(slide, 8.48, y, 4.12, 0.93, PALE_GREEN if index % 2 == 0 else CREAM, radius=True)
        add_text(slide, adjective, 8.7, y + 0.14, 3.65, 0.28, size=10.5, color=DARK_GREEN, bold=True)
        add_text(slide, context, 8.7, y + 0.5, 3.65, 0.22, size=8.5, color=MUTED, italic=True)
    add_rect(slide, 0.72, 6.43, 11.88, 0.38, GREEN)
    add_text(
        slide,
        "NATURE   ·   CULTURE   ·   RELIGION   ·   CITY",
        1.0,
        6.51,
        11.3,
        0.22,
        size=8.8,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 4)

    # Slide 5 — Activities.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, CREAM)
    add_header(slide, "04  ·  WHAT CAN PEOPLE DO?", "Activities")
    add_picture_cover(slide, PANORAMA_IMAGE, 0.72, 1.43, 3.82, 1.55)
    add_picture_cover(slide, BASILICA_IMAGE, 4.75, 1.43, 3.82, 1.55)
    add_picture_cover(slide, CABLE_CAR_IMAGE, 8.78, 1.43, 3.82, 1.55)
    activities = SLIDES[4]["screen"]
    for index, activity in enumerate(activities):
        column = index % 3
        row = index // 3
        add_numbered_card(
            slide,
            index + 1,
            activity,
            0.72 + column * 4.03,
            3.35 + row * 1.45,
            3.82,
            1.08,
            fill=WHITE,
        )
    add_footer(slide, 5)

    # Slide 6 — Opinion.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, DEEP_GREEN)
    add_header(slide, "05  ·  A PERSONAL POINT OF VIEW", "My Opinion", dark=True)
    add_picture_cover(slide, PANORAMA_IMAGE, 0.72, 1.48, 6.08, 5.35)
    add_rect(slide, 7.18, 1.48, 5.42, 5.35, CREAM, radius=True)
    add_text(slide, "“", 7.55, 1.78, 0.8, 0.72, size=45, color=GREEN, bold=True)
    add_text(
        slide,
        "In my opinion, Monserrate is one of Bogotá’s most complete attractions.",
        7.72,
        2.22,
        4.48,
        1.35,
        size=18,
        color=INK,
        bold=True,
    )
    opinion_points = [
        "My favorite feature: the panoramic view",
        "Nature + City + Tradition",
        "For families, friends, and tourists",
    ]
    for index, point in enumerate(opinion_points):
        y = 4.08 + index * 0.66
        add_rect(slide, 7.72, y + 0.04, 0.16, 0.16, GREEN, radius=True)
        add_text(slide, point, 8.08, y, 4.05, 0.35, size=10.8, color=DARK_GREEN, bold=True)
    add_footer(slide, 6, dark=True)

    # Slide 7 — Recommendations.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, PALE_GREEN)
    add_header(slide, "06  ·  BEFORE YOU GO", "Recommendations")
    recommendations = SLIDES[6]["screen"]
    for index, recommendation in enumerate(recommendations):
        column = index % 2
        row = index // 2
        add_numbered_card(
            slide,
            index + 1,
            recommendation,
            0.72 + column * 6.08,
            1.5 + row * 1.46,
            5.78,
            1.12,
            fill=WHITE if column == 0 else DARK_GREEN,
            dark=column == 1,
        )
    add_picture_cover(slide, CABLE_CAR_IMAGE, 0.72, 6.08, 11.88, 0.68)
    add_rect(slide, 8.1, 6.08, 4.5, 0.68, TERRACOTTA)
    add_text(slide, "CHECK OFFICIAL HOURS BEFORE YOU GO", 8.28, 6.27, 4.14, 0.24, size=8.3, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # Slide 8 — Conclusion and credits.
    slide = presentation.slides.add_slide(blank)
    set_slide_background(slide, DEEP_GREEN)
    add_picture_cover(slide, BASILICA_IMAGE, 6.05, 0, 7.28, 5.95)
    add_rect(slide, 5.92, 0, 0.13, 5.95, GREEN)
    add_text(slide, "CONCLUSION", 0.72, 0.7, 4.6, 0.28, size=8.8, color=SKY, bold=True)
    add_text(slide, "MONSERRATE", 0.72, 1.25, 4.8, 0.62, size=28, color=WHITE, bold=True)
    add_text(
        slide,
        "connects Bogotá with nature, culture, and tradition.",
        0.72,
        2.0,
        4.75,
        1.25,
        size=18,
        color=WHITE,
        bold=True,
    )
    add_rect(slide, 0.72, 3.72, 4.72, 0.82, GREEN, radius=True)
    add_text(slide, "THANK YOU FOR WATCHING", 0.95, 3.96, 4.25, 0.3, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Jhon Steven Alvarez Ruiz", 0.72, 5.18, 4.7, 0.3, size=10.5, color=SKY, bold=True)
    add_rect(slide, 0, 5.95, 13.333, 1.55, CREAM)
    add_text(slide, "FACTS AND VISITOR INFORMATION", 0.72, 6.12, 3.6, 0.2, size=7.4, color=DARK_GREEN, bold=True)
    add_hyperlink_line(slide, "Visit Bogotá · Cerro de Monserrate · IDRD", SOURCES[0][1], 0.72, 6.38, 4.8, size=7.2)
    add_text(slide, "IMAGE CREDITS · CC BY-SA 3.0", 6.05, 6.12, 2.9, 0.2, size=7.4, color=DARK_GREEN, bold=True)
    add_text(
        slide,
        "Jimmy Gómez N · Javadumper · diego_cue / Wikimedia Commons",
        6.05,
        6.38,
        6.55,
        0.38,
        size=7.2,
        color=MUTED,
    )
    add_text(slide, "GA2-240202501-AA2-EV02", 0.72, 7.08, 3.2, 0.18, size=6.8, color=MUTED, bold=True)
    add_text(slide, "08 / 08", 11.8, 7.08, 0.8, 0.18, size=6.8, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)

    presentation.save(PPTX_PATH)


def docx_set_cell_shading(cell, fill):
    cell_properties = cell._tc.get_or_add_tcPr()
    shading = cell_properties.find(qn("w:shd"))
    if shading is None:
        shading = OxmlElement("w:shd")
        cell_properties.append(shading)
    shading.set(qn("w:fill"), fill)


def docx_set_cell_margins(cell, top=100, start=120, bottom=100, end=120):
    cell_properties = cell._tc.get_or_add_tcPr()
    margins = cell_properties.first_child_found_in("w:tcMar")
    if margins is None:
        margins = OxmlElement("w:tcMar")
        cell_properties.append(margins)
    for name, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = margins.find(qn(f"w:{name}"))
        if node is None:
            node = OxmlElement(f"w:{name}")
            margins.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def docx_set_table_borders(table, color="D8DED9", size=5):
    table_properties = table._tbl.tblPr
    borders = table_properties.find(qn("w:tblBorders"))
    if borders is None:
        borders = OxmlElement("w:tblBorders")
        table_properties.append(borders)
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        element = OxmlElement(f"w:{edge}")
        element.set(qn("w:val"), "single")
        element.set(qn("w:sz"), str(size))
        element.set(qn("w:space"), "0")
        element.set(qn("w:color"), color)
        borders.append(element)


def docx_style_run(run, size=10.5, color=INK, bold=False, italic=False):
    run.font.name = "Lato"
    run._element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), "Lato")
    run.font.size = DocxPt(size)
    run.font.color.rgb = DocxRGBColor.from_string(color)
    run.font.bold = bold
    run.font.italic = italic


def docx_add_heading(document, text, level=1):
    paragraph = document.add_paragraph()
    paragraph.paragraph_format.space_before = DocxPt(9 if level == 1 else 6)
    paragraph.paragraph_format.space_after = DocxPt(4)
    run = paragraph.add_run(text)
    docx_style_run(run, size=15 if level == 1 else 11, color=DARK_GREEN, bold=True)
    paragraph_properties = paragraph._p.get_or_add_pPr()
    keep_next = OxmlElement("w:keepNext")
    paragraph_properties.append(keep_next)
    return paragraph


def docx_add_body(document, text, bold=False, italic=False, indent=0):
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    paragraph.paragraph_format.left_indent = DocxInches(indent)
    paragraph.paragraph_format.space_after = DocxPt(5)
    paragraph.paragraph_format.line_spacing = 1.1
    run = paragraph.add_run(text)
    docx_style_run(run, size=10.25, color=INK, bold=bold, italic=italic)
    return paragraph


def build_guide():
    document = Document()
    section = document.sections[0]
    section.page_width = DocxInches(8.5)
    section.page_height = DocxInches(11)
    section.top_margin = DocxInches(0.62)
    section.bottom_margin = DocxInches(0.62)
    section.left_margin = DocxInches(0.72)
    section.right_margin = DocxInches(0.72)

    document.core_properties.title = "Oral script and pronunciation guide — Monserrate"
    document.core_properties.subject = "GA2-240202501-AA2-EV02"
    document.core_properties.author = "Jhon Steven Alvarez Ruiz"

    normal = document.styles["Normal"]
    normal.font.name = "Lato"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Lato")
    normal.font.size = DocxPt(10.25)
    normal.font.color.rgb = DocxRGBColor.from_string(INK)

    brand = document.add_table(rows=1, cols=1)
    brand.alignment = WD_TABLE_ALIGNMENT.CENTER
    brand_cell = brand.cell(0, 0)
    docx_set_cell_shading(brand_cell, GREEN)
    docx_set_cell_margins(brand_cell, top=80, start=100, bottom=80, end=100)
    brand_p = brand_cell.paragraphs[0]
    brand_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    brand_run = brand_p.add_run("SENA  ·  ENGLISH EVIDENCE")
    docx_style_run(brand_run, size=12, color=WHITE, bold=True)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.paragraph_format.space_before = DocxPt(22)
    title_run = title.add_run("MONSERRATE\nBOGOTÁ FROM THE SKY")
    docx_style_run(title_run, size=24, color=DARK_GREEN, bold=True)

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle.add_run("Oral script, pronunciation guide, and recording checklist")
    docx_style_run(subtitle_run, size=11, color=MUTED, italic=True)

    data = document.add_table(rows=1, cols=2)
    data.alignment = WD_TABLE_ALIGNMENT.CENTER
    docx_set_table_borders(data)
    for row, (label, value) in enumerate(
        [
            ("Aprendiz", "Jhon Steven Alvarez Ruiz"),
        ]
    ):
        left, right = data.rows[row].cells
        left.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        right.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        docx_set_cell_shading(left, PALE_GREEN)
        docx_set_cell_shading(right, WHITE if row % 2 == 0 else "F5F6F4")
        docx_set_cell_margins(left)
        docx_set_cell_margins(right)
        label_run = left.paragraphs[0].add_run(label)
        value_run = right.paragraphs[0].add_run(value)
        docx_style_run(label_run, size=9.5, color=DARK_GREEN, bold=True)
        docx_style_run(value_run, size=9.8, color=INK)

    document.add_page_break()
    docx_add_heading(document, "How to use this guide")
    docx_add_body(
        document,
        "The presentation is ready, but the final video must use the learner’s own voice because "
        "the rubric evaluates pronunciation, rhythm, and intonation. Practice each slide separately, "
        "then record the complete presentation at a calm and natural pace.",
        bold=True,
    )

    timing = document.add_table(rows=1 + len(SLIDES), cols=3)
    timing.alignment = WD_TABLE_ALIGNMENT.CENTER
    docx_set_table_borders(timing)
    for column, value in enumerate(("Slide", "Topic", "Target time")):
        cell = timing.cell(0, column)
        docx_set_cell_shading(cell, DARK_GREEN)
        run = cell.paragraphs[0].add_run(value)
        docx_style_run(run, size=9, color=WHITE, bold=True)
    for row, slide_info in enumerate(SLIDES, start=1):
        for column, value in enumerate((str(row), slide_info["title"], slide_info["time"])):
            cell = timing.cell(row, column)
            docx_set_cell_shading(cell, WHITE if row % 2 else "F5F6F4")
            run = cell.paragraphs[0].add_run(value)
            docx_style_run(run, size=8.8, color=INK, bold=column == 0)

    for index, slide_info in enumerate(SLIDES, start=1):
        if index in (1, 3, 5, 7):
            document.add_page_break()
        docx_add_heading(document, f"Slide {index} — {slide_info['title']}")
        time_p = document.add_paragraph()
        time_run = time_p.add_run(f"Target: {slide_info['time']}  ·  Pronunciation focus: {slide_info['focus']}")
        docx_style_run(time_run, size=8.8, color=GREEN, bold=True)

        screen_p = document.add_paragraph()
        screen_run = screen_p.add_run("ON-SCREEN IDEAS")
        docx_style_run(screen_run, size=8.2, color=MUTED, bold=True)
        for item in slide_info["screen"]:
            bullet = document.add_paragraph(style=None)
            bullet.paragraph_format.left_indent = DocxInches(0.24)
            bullet.paragraph_format.space_after = DocxPt(1.5)
            bullet_run = bullet.add_run(f"• {item}")
            docx_style_run(bullet_run, size=9.3, color=INK)

        script_p = document.add_paragraph()
        script_run = script_p.add_run("SPEAKING SCRIPT")
        docx_style_run(script_run, size=8.2, color=MUTED, bold=True)
        quote = document.add_table(rows=1, cols=1)
        quote.alignment = WD_TABLE_ALIGNMENT.CENTER
        quote_cell = quote.cell(0, 0)
        docx_set_cell_shading(quote_cell, PALE_GREEN)
        docx_set_cell_margins(quote_cell, top=130, start=160, bottom=130, end=160)
        quote_run = quote_cell.paragraphs[0].add_run(slide_info["script"])
        docx_style_run(quote_run, size=10.1, color=INK)

    document.add_page_break()
    docx_add_heading(document, "Pronunciation reference")
    docx_add_body(
        document,
        "The last column is only an approximate aid for Spanish speakers. Listen to the stress shown "
        "in capital letters, pronounce final consonants, and avoid reading too quickly.",
        italic=True,
    )
    pronunciation_table = document.add_table(rows=1 + len(PRONUNCIATION), cols=3)
    pronunciation_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    docx_set_table_borders(pronunciation_table)
    for column, value in enumerate(("Word", "IPA", "Approximate support")):
        cell = pronunciation_table.cell(0, column)
        docx_set_cell_shading(cell, DARK_GREEN)
        run = cell.paragraphs[0].add_run(value)
        docx_style_run(run, size=8.8, color=WHITE, bold=True)
    for row, values in enumerate(PRONUNCIATION, start=1):
        for column, value in enumerate(values):
            cell = pronunciation_table.cell(row, column)
            docx_set_cell_shading(cell, WHITE if row % 2 else "F5F6F4")
            run = cell.paragraphs[0].add_run(value)
            docx_style_run(run, size=8.2, color=INK, bold=column == 0)

    document.add_page_break()
    docx_add_heading(document, "Recording checklist")
    checklist = [
        "Use your own voice; do not submit synthetic narration.",
        "Record in landscape format and show the slides in full screen.",
        "Choose a quiet room and test the microphone before the final take.",
        "Speak at a calm pace, pause between ideas, and vary your intonation.",
        "Look at the camera during the introduction and conclusion if you appear on screen.",
        "Use the script as support, but do not read every slide bullet word for word.",
        "Check that the final video opens correctly and that the audio is clear.",
        "Upload the file or permitted link to the official learning platform before the deadline.",
    ]
    for item in checklist:
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.left_indent = DocxInches(0.15)
        paragraph.paragraph_format.space_after = DocxPt(5)
        run = paragraph.add_run(f"☐  {item}")
        docx_style_run(run, size=10, color=INK)

    docx_add_heading(document, "Sources and image credits")
    for label, url in SOURCES + IMAGE_CREDITS:
        paragraph = document.add_paragraph()
        paragraph.paragraph_format.left_indent = DocxInches(0.16)
        paragraph.paragraph_format.space_after = DocxPt(3)
        run = paragraph.add_run(f"• {label}\n  {url}")
        docx_style_run(run, size=8.3, color=DARK_GREEN)

    document.save(GUIDE_PATH)


def build_markdown():
    lines = [
        "# Monserrate: Bogotá from the Sky",
        "",
        "**Evidencia:** GA2-240202501-AA2-EV02",
        "",
        "**Aprendiz:** Jhon Steven Alvarez Ruiz",
        "",
        "> El video final debe grabarse con la voz del aprendiz, porque la rúbrica evalúa pronunciación, ritmo y entonación.",
        "",
    ]
    for index, slide_info in enumerate(SLIDES, start=1):
        lines.extend(
            [
                f"## Slide {index} — {slide_info['title']}",
                "",
                f"**Tiempo objetivo:** {slide_info['time']}",
                "",
                "**Ideas visibles:**",
                "",
            ]
        )
        lines.extend(f"- {item}" for item in slide_info["screen"])
        lines.extend(
            [
                "",
                "**Guion oral:**",
                "",
                slide_info["script"],
                "",
                f"**Pronunciación a practicar:** {slide_info['focus']}",
                "",
            ]
        )
    lines.extend(["## Fuentes", ""])
    for label, url in SOURCES:
        lines.append(f"- [{label}]({url})")
    lines.extend(["", "## Créditos de imágenes", ""])
    for label, url in IMAGE_CREDITS:
        lines.append(f"- [{label}]({url})")
    MARKDOWN_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
    build_guide()
    build_markdown()
    total_words = sum(len(slide["script"].split()) for slide in SLIDES)
    print(f"Created: {PPTX_PATH}")
    print(f"Created: {GUIDE_PATH}")
    print(f"Created: {MARKDOWN_PATH}")
    print(f"Oral script word count: {total_words}")
